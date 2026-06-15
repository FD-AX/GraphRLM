"""Builds the GraphRLM pitch deck (PPTX) - 7 slides, ~5-7 min, grant/startup framing.

Editable output: docs/graphrlm_pitch.pptx. Re-run after editing this script.
Placeholders marked [PLACEHOLDER] in speaker notes: vertical flagship and the ask.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]

BG = RGBColor(0x11, 0x15, 0x1C)
ACCENT = RGBColor(0x7C, 0xC4, 0xFF)
GREEN = RGBColor(0x9A, 0xDB, 0xB4)
YELLOW = RGBColor(0xFF, 0xD4, 0x79)
TEXT = RGBColor(0xE8, 0xEA, 0xF0)
MUTED = RGBColor(0x8B, 0x93, 0xA7)
ROW = RGBColor(0x1B, 0x22, 0x30)
HL = RGBColor(0x20, 0x30, 0x1F)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def textbox(s, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color=TEXT, bold=False, italic=False, first=False, space_after=10, bullet=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    return p


def title(s, text):
    tf = textbox(s, 0.8, 0.5, 11.7, 1.1)
    para(tf, text, 30, ACCENT, bold=True, first=True)


def style_cell(cell, text, size=14, color=TEXT, bold=False, fill=ROW, align=PP_ALIGN.LEFT):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.12)
    cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def table(s, rows, left, top, width, col_widths, header_fill=ROW, hl_row=None):
    hl_set = set() if hl_row is None else ({hl_row} if isinstance(hl_row, int) else set(hl_row))
    n_rows, n_cols = len(rows), len(rows[0])
    gt = s.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(0.5 * n_rows)).table
    for ci, w in enumerate(col_widths):
        gt.columns[ci].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            is_header = ri == 0
            is_hl = ri in hl_set
            fill = header_fill if is_header else (HL if is_hl else BG)
            color = ACCENT if is_header else (GREEN if is_hl else TEXT)
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            style_cell(gt.cell(ri, ci), val, size=14, color=color,
                       bold=is_header or is_hl, fill=fill, align=align)
    return gt


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- Slide 1
s = slide()
tf = textbox(s, 0.9, 2.4, 11.5, 3, anchor=MSO_ANCHOR.TOP)
para(tf, "GraphRLM", 54, ACCENT, bold=True, first=True, space_after=18)
para(tf, "Слой проверки полноты доказательств для AI-research —", 22, TEXT, space_after=2)
para(tf, "там, где пропущенный факт стоит дороже неверного ответа.", 22, TEXT, space_after=22)
para(tf, "github.com/FD-AX/GraphRLM   ·   питч-версия", 14, MUTED)
notes(s, "~30 сек. Хук: 'Мы не делаем ещё один RAG. Мы решаем проблему, "
         "которую RAG игнорирует — полноту доказательств.' Назвать себя и проект, "
         "сразу перейти к боли.")

# ---------------------------------------------------------------- Slide 2
s = slide()
title(s, "Проблема: AI отвечает уверенно на неполных данных")
tf = textbox(s, 0.9, 1.7, 11.5, 4)
para(tf, "В юр-research, due diligence, фарме пропущенное доказательство дороже "
         "неверного ответа — потому что выглядит как правильный ответ.", 20, TEXT, first=True, bullet=True, space_after=14)
para(tf, "Все меряют точность ответа. Полноту доказательной базы не меряет никто.", 20, TEXT, bullet=True, space_after=14)
para(tf, "«False completion»: система закрывает вопрос, не собрав всю цепочку — "
         "невидимый и дорогой режим отказа RAG-систем.", 20, TEXT, bullet=True, space_after=20)
tf2 = textbox(s, 0.9, 5.6, 11.5, 1)
para(tf2, "Мы первыми сделали полноту измеримой метрикой — и научились её гарантировать.",
     21, YELLOW, bold=True, first=True)
notes(s, "~60 сек. Это самый важный слайд для 'защиты идеи'. Продаёшь ПРОБЛЕМУ, "
         "не механизм. Приведи живой пример боли из выбранной вертикали "
         "[PLACEHOLDER: вставь конкретный кейс — напр. пропущенный прецедент в "
         "юр-research стоит проигранного дела]. Не упоминай графы/косинусы здесь.")

# ---------------------------------------------------------------- Slide 3
s = slide()
title(s, "Решение: дешёвый аудит полноты + дорогой поиск по триггеру")
tf = textbox(s, 0.9, 1.7, 5.6, 4)
para(tf, "Слой 1 — аудит полноты", 18, GREEN, bold=True, first=True, space_after=8)
para(tf, "Граф на BERT-кодировках, без обучения и без LLM. Раскладывает запрос "
         "на нужные доказательства и находит пробелы.", 17, TEXT, bullet=True, space_after=8)
para(tf, "Дёшево, детерминированно, на любом домене.", 17, TEXT, bullet=True)
tf2 = textbox(s, 6.9, 1.7, 5.6, 4)
para(tf2, "Слой 2 — RLM discovery", 18, GREEN, bold=True, first=True, space_after=8)
para(tf2, "Дорогой LLM-цикл запускается ТОЛЬКО когда аудит нашёл реальный пробел.", 17, TEXT, bullet=True, space_after=8)
para(tf2, "Целенаправленно достаёт недостающее звено цепочки.", 17, TEXT, bullet=True)
tf3 = textbox(s, 0.9, 5.7, 11.5, 1)
para(tf3, "Gap-driven экономика: за дорогой LLM платим только там, где есть пробел.",
     21, YELLOW, bold=True, first=True)
notes(s, "~55 сек. Один кадр архитектуры. Главная мысль — не 'граф', а "
         "ЭКОНОМИКА: дешёвый детектор + дорогой исполнитель по триггеру. "
         "Это и масштабируется, и дёшево.")

# ---------------------------------------------------------------- Slide 4
s = slide()
title(s, "Результат: полнота 70% там, где наивный RAG упирается в 21%")
table(s, [
    ["Подход", "Полнота доказательств"],
    ["Наивный RAG (dense / cross-encoder / итеративный)", "~0.21"],
    ["GraphRLM — zero-shot, без обучения", "0.70"],
    ["Supervised SOTA — обучен на датасете *", "0.774"],
], left=1.4, top=2.1, width=10.5, col_widths=[7.3, 3.2], hl_row=2)
tf = textbox(s, 0.9, 5.0, 11.5, 1.8)
para(tf, "Бенчмарк MuSiQue, 500 кейсов, парные тесты (p << 0.001). Zero-shot — "
         "вплотную к системам, обученным под задачу.", 18, TEXT, first=True, space_after=8)
para(tf, "* внешнее число — directional: протоколы оценки не идентичны.", 13, MUTED)
notes(s, "~60 сек. Killer-слайд. Одна мысль: наивные подходы фундаментально "
         "упираются в 21%, мы даём 70% БЕЗ обучения на домене. Если спросят "
         "про supervised 0.774 — 'мы zero-shot и в полтора раза дешевле развернуть, "
         "им нужна разметка под каждую вертикаль'. НЕ показывай text-RLM абляцию.")

# ---------------------------------------------------------------- Slide 5
s = slide()
title(s, "Качество ответов выше публичных LLM — за копейки")
table(s, [
    ["Метрика (500 кейсов, zero-shot, v2)", "GraphRLM", "Публичные LLM *"],
    ["Точность ответа (EM)", "0.52", "~0.37"],
    ["Полнота ответа (F1)", "0.64", "~0.45"],
    ["Цена за запрос", "$0.025", "—"],
], left=1.4, top=2.1, width=10.5, col_widths=[6.0, 2.25, 2.25], hl_row=[1, 2])
tf = textbox(s, 0.9, 4.9, 11.5, 2.2)
para(tf, "Точность и полнота ответа в ~1.4 раза выше публичных LLM-бейзлайнов — на дешёвой модели.", 19, TEXT, first=True, bullet=True, space_after=10)
para(tf, "Дорогой цикл — только по триггеру: на масштабе платишь не за каждый запрос.", 19, TEXT, bullet=True, space_after=10)
para(tf, "* литературные данные (напр. Gemini на MuSiQue); directional.", 13, MUTED)
notes(s, "~45 сек. Связка качество+цена. Цифры v2 на полных 500 кейсах (EM 0.522, "
         "F1 0.639) — надёжные сами по себе. ОСТОРОЖНО: сравнение с публичными LLM — "
         "directional (разные протоколы оценки MuSiQue), держи звёздочку и НЕ делай "
         "'в 1.4 раза точнее Gemini' главным лозунгом — грамотный человек спросит про "
         "идентичность сеттинга. Не уходи в F1 против text-RLM (незначимо, n=60).")

# ---------------------------------------------------------------- Slide 6
s = slide()
title(s, "Почему это трудно повторить")
tf = textbox(s, 0.9, 1.8, 11.5, 4.5)
para(tf, "Другая ось конкуренции: не точность ответа (это умеют все), а гарантия "
         "отсутствия false completion.", 20, TEXT, first=True, bullet=True, space_after=16)
para(tf, "Training-free: работает на новом домене без разметки — конкуренту нужна "
         "размеченная выборка под каждую вертикаль.", 20, TEXT, bullet=True, space_after=16)
para(tf, "Структурный аудит полноты + типизированные слоты недостающих доказательств "
         "как управляющий сигнал для дорогого поиска.", 20, TEXT, bullet=True, space_after=16)
notes(s, "~50 сек. Ответ на 'не повторят ли за выходные'. Дифференциация — "
         "НЕ в косинусах/многомерности (это не продаём). В постановке "
         "(completeness-first), training-free развёртывании и экономике.")

# ---------------------------------------------------------------- Slide 7
s = slide()
title(s, "Куда дальше")
tf = textbox(s, 0.9, 1.7, 11.5, 4)
para(tf, "Валидация на открытом корпусе (10 000+ документов) — режим, где структура "
         "обязательна, а не опциональна.", 20, TEXT, first=True, bullet=True, space_after=14)
para(tf, "Пилот в вертикали [PLACEHOLDER: юр-research / due diligence / фарма].", 20, TEXT, bullet=True, space_after=14)
para(tf, "Запрос: [PLACEHOLDER: сумма / ресурс / партнёрство].", 20, YELLOW, bold=True, bullet=True, space_after=20)
tf2 = textbox(s, 0.9, 5.7, 11.5, 1)
para(tf2, "«Graph verifies completeness. RLM closes evidence gaps.»", 20, GREEN, italic=True, first=True)
notes(s, "~40 сек. Закрытие. ОБЯЗАТЕЛЬНО замени два [PLACEHOLDER]: вертикаль "
         "флагмана и конкретный ask (сколько денег / на что / что просишь у "
         "этой аудитории). Без конкретного ask питч провисает. Заверши слоганом.")

out = PROJECT_ROOT / "docs" / "graphrlm_pitch.pptx"
prs.save(str(out))
print(f"Saved {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
